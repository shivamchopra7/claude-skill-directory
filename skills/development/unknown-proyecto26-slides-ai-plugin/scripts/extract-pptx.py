#!/usr/bin/env python3
"""
Extract text, images, and speaker notes from a PPTX file to JSON.
Used for converting existing PowerPoint presentations to HTML slides.

Usage:
    python extract-pptx.py <input.pptx> [--output extracted-slides.json] [--assets-dir assets/]

Dependencies:
    pip install python-pptx Pillow --break-system-packages
"""

import argparse
import json
import os
import sys

def extract_pptx(pptx_path, output_path="extracted-slides.json", assets_dir="assets/"):
    """Extract slides from a PPTX file into structured JSON."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
    except ImportError:
        print("Error: python-pptx is required. Install with:")
        print("  pip install python-pptx --break-system-packages")
        sys.exit(1)

    if not os.path.exists(pptx_path):
        print(f"Error: File not found: {pptx_path}")
        sys.exit(1)

    os.makedirs(assets_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    slides_data = []
    image_count = 0

    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": slide_num,
            "title": "",
            "content": [],
            "images": [],
            "notes": "",
            "layout": "content"
        }

        # Extract shapes
        for shape in slide.shapes:
            # Title
            if shape.has_text_frame and shape.shape_id == slide.shapes.title.shape_id if slide.shapes.title else False:
                slide_info["title"] = shape.text_frame.text.strip()
                continue

            # Text content
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_info["content"].append({
                            "text": text,
                            "level": paragraph.level,
                            "bold": any(run.font.bold for run in paragraph.runs if run.font.bold),
                            "font_size": str(paragraph.runs[0].font.size) if paragraph.runs and paragraph.runs[0].font.size else None
                        })

            # Images
            if shape.shape_type == 13:  # Picture type
                image_count += 1
                image_name = f"slide{slide_num}_img{image_count}.png"
                image_path = os.path.join(assets_dir, image_name)

                try:
                    image_blob = shape.image.blob
                    with open(image_path, "wb") as f:
                        f.write(image_blob)
                    slide_info["images"].append({
                        "filename": image_name,
                        "path": image_path,
                        "width": shape.width,
                        "height": shape.height
                    })
                except Exception as e:
                    print(f"  Warning: Could not extract image from slide {slide_num}: {e}")

        # If no title found via shape_id, use first text shape
        if not slide_info["title"] and slide_info["content"]:
            slide_info["title"] = slide_info["content"][0]["text"]
            slide_info["content"] = slide_info["content"][1:]

        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            slide_info["notes"] = slide.notes_slide.notes_text_frame.text.strip()

        # Detect layout type
        if not slide_info["content"] and slide_info["images"]:
            slide_info["layout"] = "image-focus"
        elif slide_info["images"] and slide_info["content"]:
            slide_info["layout"] = "two-column"
        elif slide_num == 1:
            slide_info["layout"] = "title"
        elif len(slide_info["content"]) == 0:
            slide_info["layout"] = "section-divider"

        slides_data.append(slide_info)

    # Write output
    output = {
        "source": os.path.basename(pptx_path),
        "total_slides": len(slides_data),
        "slides": slides_data
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(output, f, indent=2, ensure_ascii=False, default=str)

    print(f"Extracted {len(slides_data)} slides to {output_path}")
    print(f"Saved {image_count} images to {assets_dir}")
    return output


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract PPTX slides to JSON for HTML conversion")
    parser.add_argument("pptx_path", help="Path to the .pptx file")
    parser.add_argument("--output", "-o", default="extracted-slides.json", help="Output JSON path")
    parser.add_argument("--assets-dir", "-a", default="assets/", help="Directory for extracted images")
    args = parser.parse_args()

    extract_pptx(args.pptx_path, args.output, args.assets_dir)
